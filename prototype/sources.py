from __future__ import annotations
from io import BytesIO
from pathlib import Path
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

MAX_BYTES = 25 * 1024 * 1024
ALLOWED = {"application/pdf": ".pdf", "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx", "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx", "text/plain": ".txt"}

def validate_upload(filename: str, content_type: str, content: bytes) -> None:
    if not filename or Path(filename).suffix.lower() not in ALLOWED.values() or content_type not in ALLOWED:
        raise ValueError("Only PDF, DOCX, PPTX and text files are supported.")
    if not content or len(content) > MAX_BYTES: raise ValueError("File must be between 1 byte and 25 MB.")

def extract_text(content_type: str, content: bytes) -> str:
    if content_type == "text/plain": return content.decode("utf-8", errors="replace")
    if content_type == "application/pdf": return "\n".join(page.extract_text() or "" for page in PdfReader(BytesIO(content)).pages)
    if content_type.endswith("wordprocessingml.document"): return "\n".join(p.text for p in Document(BytesIO(content)).paragraphs)
    if content_type.endswith("presentationml.presentation"):
        return "\n".join(shape.text for slide in Presentation(BytesIO(content)).slides for shape in slide.shapes if hasattr(shape, "text"))
    return ""
